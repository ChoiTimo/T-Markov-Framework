"""PPTX assembler — Phase 2 Sprint 2-3.

Builds a PowerPoint deck from slide instance dicts.
Keeps rendering simple and self-contained (no external templates required).

Font note: python-pptx uses Cairo-style embedding-less text; rendering
relies on the viewer's font. Korean should render correctly on any
PowerPoint install with Malgun Gothic / Apple SD Gothic Neo available.
"""

from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# ------------------------------------------------------------------
# Theme
# ------------------------------------------------------------------

SLIDE_WIDTH_EMU = Inches(13.333)   # 16:9 widescreen
SLIDE_HEIGHT_EMU = Inches(7.5)

COLOR_PRIMARY    = RGBColor(0x0F, 0x1F, 0x3D)   # dark navy
COLOR_ACCENT     = RGBColor(0xE5, 0x00, 0x3C)   # crimson
COLOR_TEXT       = RGBColor(0x1A, 0x1A, 0x1A)
COLOR_MUTED      = RGBColor(0x66, 0x66, 0x66)
COLOR_LIGHT_BG   = RGBColor(0xF5, 0xF6, 0xF8)

# Phase 별 헤더 배지 색
PHASE_COLORS = {
    "frame":      RGBColor(0x20, 0x4E, 0x8A),
    "tension":    RGBColor(0xB0, 0x50, 0x1B),
    "surprise":   RGBColor(0xA0, 0x1A, 0x58),
    "evidence":   RGBColor(0x1F, 0x6E, 0x4A),
    "conviction": RGBColor(0x0F, 0x1F, 0x3D),
}

PHASE_LABELS = {
    "frame":      "Phase 1 · Frame",
    "tension":    "Phase 2 · Tension",
    "surprise":   "Phase 3 · Surprise",
    "evidence":   "Phase 4 · Evidence",
    "conviction": "Phase 5 · Conviction",
}


# ------------------------------------------------------------------
# Low-level helpers
# ------------------------------------------------------------------

def _add_blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _add_textbox(slide, left, top, width, height, text, *,
                 font_size=18, bold=False, color=COLOR_TEXT,
                 align=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def _add_phase_header(slide, phase: str):
    label = PHASE_LABELS.get(phase, "")
    color = PHASE_COLORS.get(phase, COLOR_PRIMARY)
    if not label:
        return
    _add_textbox(
        slide,
        left=Inches(0.6), top=Inches(0.35),
        width=Inches(5), height=Inches(0.4),
        text=label, font_size=11, bold=True, color=color,
    )


def _add_title(slide, text: str, *, color=COLOR_PRIMARY):
    _add_textbox(
        slide,
        left=Inches(0.6), top=Inches(0.85),
        width=Inches(12.1), height=Inches(1.0),
        text=text or "", font_size=30, bold=True, color=color,
    )


def _add_subtitle(slide, text: str | None):
    if not text:
        return
    _add_textbox(
        slide,
        left=Inches(0.6), top=Inches(1.85),
        width=Inches(12.1), height=Inches(0.5),
        text=text, font_size=16, color=COLOR_MUTED,
    )


def _add_bullet_list(slide, bullets: list[str], *, top=Inches(2.5), font_size=20):
    tb = slide.shapes.add_textbox(
        Inches(0.6), top, Inches(12.1), Inches(4.5)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets or []):
        if not text:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {text}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = COLOR_TEXT
        run.font.name = "맑은 고딕"


def _add_footer(slide, proposal_number: str | None, page_no: int, total: int):
    footer_text = f"{proposal_number or ''}   {page_no} / {total}"
    _add_textbox(
        slide,
        left=Inches(0.6), top=Inches(7.05),
        width=Inches(12.1), height=Inches(0.3),
        text=footer_text, font_size=10, color=COLOR_MUTED,
    )


# ------------------------------------------------------------------
# Per-phase renderers
# ------------------------------------------------------------------

def _render_cover(slide, proposal: dict, body: dict):
    customer = proposal.get("customer_company") or proposal.get("customer_name") or ""
    title = body.get("title") or proposal.get("title") or "SmartWAN 제안서"
    subtitle = body.get("subtitle") or (customer and f"{customer} 귀중")
    date = body.get("date") or datetime.utcnow().strftime("%Y.%m.%d")
    author = body.get("author_name") or ""

    _add_textbox(
        slide,
        left=Inches(0.9), top=Inches(2.2),
        width=Inches(11.5), height=Inches(1.4),
        text=title, font_size=40, bold=True, color=COLOR_PRIMARY,
    )
    if subtitle:
        _add_textbox(
            slide,
            left=Inches(0.9), top=Inches(3.8),
            width=Inches(11.5), height=Inches(0.8),
            text=subtitle, font_size=22, color=COLOR_TEXT,
        )
    meta_line = " | ".join([x for x in (date, author) if x])
    if meta_line:
        _add_textbox(
            slide,
            left=Inches(0.9), top=Inches(6.6),
            width=Inches(11.5), height=Inches(0.4),
            text=meta_line, font_size=12, color=COLOR_MUTED,
        )


def _render_narrative(slide, body: dict):
    headline = body.get("headline") or "하나의 스토리"
    _add_textbox(
        slide,
        left=Inches(0.6), top=Inches(0.85),
        width=Inches(12.1), height=Inches(1.0),
        text=headline, font_size=28, bold=True, color=COLOR_PRIMARY,
    )
    narrative = body.get("narrative") or ""
    _add_textbox(
        slide,
        left=Inches(0.9), top=Inches(2.6),
        width=Inches(11.5), height=Inches(2.5),
        text=narrative, font_size=24, color=COLOR_TEXT,
    )
    support = body.get("supporting_insight")
    if support:
        _add_textbox(
            slide,
            left=Inches(0.9), top=Inches(5.3),
            width=Inches(11.5), height=Inches(1.2),
            text=f"핵심 인사이트: {support}",
            font_size=16, color=COLOR_ACCENT, bold=True,
        )


def _render_matrix(slide, body: dict):
    headline = body.get("headline") or ""
    matrix = body.get("matrix") or []
    _add_title(slide, headline)

    rows = len(matrix) + 1
    cols = 3
    left = Inches(0.6)
    top = Inches(2.5)
    width = Inches(12.1)
    height = Inches(0.6 * rows)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    headers = ["항목", "제안", "기존 방식"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.color.rgb = COLOR_PRIMARY

    for r, row in enumerate(matrix, start=1):
        cells = [row.get("capability", ""), row.get("ours", ""), row.get("legacy", "")]
        for c, text in enumerate(cells):
            cell = table.cell(r, c)
            cell.text = str(text)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)


def _render_metrics(slide, body: dict):
    headline = body.get("headline") or ""
    metrics = body.get("metrics") or []
    note = body.get("note") or ""
    _add_title(slide, headline)

    rows = len(metrics) + 1
    cols = 3
    table_shape = slide.shapes.add_table(
        rows, cols, Inches(0.6), Inches(2.5), Inches(12.1), Inches(0.6 * rows)
    )
    table = table_shape.table
    for c, h in enumerate(["지표", "기존", "제안 후"]):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(14)
    for r, row in enumerate(metrics, start=1):
        for c, text in enumerate([row.get("label", ""), row.get("before", ""), row.get("after", "")]):
            cell = table.cell(r, c)
            cell.text = str(text)
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
    if note:
        _add_textbox(
            slide,
            left=Inches(0.6), top=Inches(6.0),
            width=Inches(12.1), height=Inches(0.5),
            text=note, font_size=11, color=COLOR_MUTED,
        )


def _render_roi(slide, body: dict):
    _add_title(slide, body.get("headline") or "ROI / TCO")
    total = body.get("total_amount")
    monthly = body.get("monthly_amount")
    contract_months = body.get("contract_months")

    top = Inches(2.4)
    if monthly:
        _add_textbox(
            slide,
            left=Inches(0.6), top=top,
            width=Inches(6), height=Inches(0.6),
            text=f"월 이용료(세후): ₩{int(float(monthly)):,}",
            font_size=22, bold=True, color=COLOR_ACCENT,
        )
        top = Inches(3.1)
    if total:
        _add_textbox(
            slide,
            left=Inches(0.6), top=top,
            width=Inches(12.1), height=Inches(0.5),
            text=f"약정 {contract_months}개월 총액(세후): ₩{int(float(total)):,}",
            font_size=16, color=COLOR_TEXT,
        )

    _add_bullet_list(slide, body.get("bullets") or [], top=Inches(4.0), font_size=16)

    disclaimer = body.get("disclaimer") or ""
    if disclaimer:
        _add_textbox(
            slide,
            left=Inches(0.6), top=Inches(6.2),
            width=Inches(12.1), height=Inches(0.5),
            text=disclaimer, font_size=10, color=COLOR_MUTED,
        )


def _render_roadmap(slide, body: dict):
    _add_title(slide, body.get("headline") or "도입 로드맵")
    phases = body.get("phases") or []
    bullets = [
        f"{ph.get('name', '')} ({ph.get('weeks', '')}주): {ph.get('desc', '')}"
        for ph in phases
    ]
    _add_bullet_list(slide, bullets, top=Inches(2.6), font_size=18)


def _render_cta(slide, body: dict):
    _add_title(slide, body.get("headline") or "Next Step", color=COLOR_ACCENT)
    bullets = body.get("bullets") or []
    _add_bullet_list(slide, bullets, top=Inches(2.6), font_size=22)
    cta = body.get("call_to_action") or ""
    if cta:
        _add_textbox(
            slide,
            left=Inches(0.6), top=Inches(6.0),
            width=Inches(12.1), height=Inches(0.8),
            text=cta, font_size=14, color=COLOR_PRIMARY, bold=True,
        )


def _render_question(slide, body: dict):
    headline = body.get("headline") or ""
    subtext = body.get("subtext") or ""
    _add_textbox(
        slide,
        left=Inches(0.6), top=Inches(2.8),
        width=Inches(12.1), height=Inches(1.6),
        text=headline, font_size=36, bold=True, color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    if subtext:
        _add_textbox(
            slide,
            left=Inches(0.6), top=Inches(4.6),
            width=Inches(12.1), height=Inches(1.0),
            text=subtext, font_size=18, color=COLOR_MUTED,
            align=PP_ALIGN.CENTER,
        )


def _render_surprise(slide, body: dict):
    _add_title(slide, body.get("headline") or "예상과 다른 지점", color=COLOR_ACCENT)
    twist = body.get("twist") or ""
    _add_textbox(
        slide,
        left=Inches(0.9), top=Inches(2.7),
        width=Inches(11.5), height=Inches(2.5),
        text=twist, font_size=22, color=COLOR_TEXT, bold=True,
    )
    hint = body.get("evidence_hint") or ""
    if hint:
        _add_textbox(
            slide,
            left=Inches(0.9), top=Inches(5.6),
            width=Inches(11.5), height=Inches(0.8),
            text=hint, font_size=14, color=COLOR_MUTED,
        )


def _render_generic(slide, body: dict):
    _add_title(slide, body.get("headline") or "")
    if body.get("subtitle"):
        _add_subtitle(slide, body.get("subtitle"))
    if body.get("bullets"):
        _add_bullet_list(slide, body.get("bullets") or [])
    elif body.get("scenario"):
        _add_textbox(
            slide,
            left=Inches(0.9), top=Inches(2.6),
            width=Inches(11.5), height=Inches(4),
            text=body.get("scenario"), font_size=18, color=COLOR_TEXT,
        )


# Builders by code
_SLIDE_RENDERERS = {
    "P1_cover":          lambda s, p, b: _render_cover(s, p, b),
    "N1_narrative":      lambda s, p, b: _render_narrative(s, b),
    "N4_surprise":       lambda s, p, b: _render_surprise(s, b),
    "N3_question_risk":  lambda s, p, b: _render_question(s, b),
    "N5_question_proof": lambda s, p, b: _render_question(s, b),
    "N6_call_to_action": lambda s, p, b: _render_cta(s, b),
    "P9_capability":     lambda s, p, b: _render_matrix(s, b),
    "P10_performance":   lambda s, p, b: _render_metrics(s, b),
    "P11_roi":           lambda s, p, b: _render_roi(s, b),
    "P12_roadmap":       lambda s, p, b: _render_roadmap(s, b),
}


def _render_slide(slide, proposal: dict, slide_data: dict):
    code = slide_data.get("code") or ""
    body = slide_data.get("body") or {}
    phase = slide_data.get("phase") or ""

    # Cover 슬라이드는 phase 헤더 없이 풀블리드로
    if code != "P1_cover":
        _add_phase_header(slide, phase)

    renderer = _SLIDE_RENDERERS.get(code)
    if renderer:
        renderer(slide, proposal, body)
    else:
        _render_generic(slide, body)


# ------------------------------------------------------------------
# Public API
# ------------------------------------------------------------------

def assemble_pptx(proposal: dict, slides: list[dict]) -> bytes:
    """Build a PPTX binary from slide instances and return bytes.

    Args:
      proposal: proposals row dict
      slides:   ordered list of slide instance dicts (output of
                neuro_optimizer.build_slide_instances)
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH_EMU
    prs.slide_height = SLIDE_HEIGHT_EMU

    total = len(slides)
    proposal_number = proposal.get("proposal_number") or ""

    for idx, slide_data in enumerate(slides, start=1):
        slide = _add_blank_slide(prs)
        _render_slide(slide, proposal, slide_data)
        # 표지는 페이지 번호 생략
        if slide_data.get("code") != "P1_cover":
            _add_footer(slide, proposal_number, idx, total)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
